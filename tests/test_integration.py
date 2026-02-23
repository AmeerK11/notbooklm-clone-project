"""Integration tests for the full ingestion pipeline: upload → extract → chunk → embed → store."""
import sys
import pathlib
import json
import tempfile
import shutil
import pytest
from pathlib import Path
from unittest.mock import patch, MagicMock

# Ensure `src` is on sys.path
ROOT = pathlib.Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "src"))

import ingestion.storage as storage
import ingestion.extractors as extractors
import ingestion.chunker as chunker
import ingestion.embeddings as embeddings
import ingestion.vectorstore as vectorstore


def test_txt_upload_extract_ingest(tmp_path):
    """Test end-to-end TXT upload → extract → chunk → embed → store."""
    # Create a test text file
    test_file = tmp_path / "test.txt"
    test_content = "Sentence one. Sentence two. Sentence three. " * 20
    test_file.write_text(test_content, encoding="utf-8")

    # Initialize adapter
    storage_dir = tmp_path / "storage"
    adapter = storage.LocalStorageAdapter(base_dir=str(storage_dir))

    # Step 1: Upload
    source_id = "test-source-001"
    user, notebook = "testuser", "test-notebook"
    dest = adapter.save_raw_file(user, notebook, source_id, test_file)
    assert dest.exists()

    # Step 2: Extract
    result = extractors.extract_text_from_txt(test_file)
    assert result["text"] == test_content
    assert result["pages"] == 1

    # Save extracted text
    adapter.save_extracted_text(user, notebook, source_id, "content", result["text"])
    extracted_path = storage_dir / "users" / user / "notebooks" / notebook / "files_extracted" / source_id / "content.txt"
    assert extracted_path.exists()
    assert extracted_path.read_text(encoding="utf-8") == test_content

    # Step 3: Chunk
    class DummyTokenizer:
        def encode(self, s, add_special_tokens=False):
            return [0] * max(1, len(s.split()))

    with patch.object(chunker, "get_tokenizer", lambda model_name=None: DummyTokenizer()):
        chunks = chunker.chunk_text(result["text"], model_name="dummy", chunk_size_tokens=50)
    assert len(chunks) > 1
    assert all("chunk_id" in c and "text" in c for c in chunks)

    # Attach metadata
    for c in chunks:
        c["source_id"] = source_id

    # Step 4: Embed (mock embedding to avoid model download)
    mock_embedder = MagicMock()
    mock_embeddings = [[0.1 * i for _ in range(384)] for i in range(len(chunks))]
    mock_embedder.embed_texts.return_value = mock_embeddings

    # Step 5: Store in Chroma
    chroma_dir = str((storage_dir / user / notebook / "chroma").resolve())
    store = vectorstore.ChromaAdapter(persist_directory=chroma_dir)
    store.upsert_chunks(user, notebook, chunks, mock_embeddings)

    # Verify storage
    collection = store.get_or_create_collection(user, notebook)
    assert collection.count() == len(chunks)


def test_url_extraction_with_fallback(tmp_path):
    """Test URL extraction with mocked response."""
    # Mock response
    mock_html = """
    <html>
        <body>
            <article>
                <p>This is the main content of the article.</p>
                <p>It should be extracted correctly.</p>
            </article>
            <footer>Footer text (should be filtered out)</footer>
        </body>
    </html>
    """

    with patch("ingestion.extractors.requests.get") as mock_get:
        mock_response = MagicMock()
        mock_response.text = mock_html
        mock_response.raise_for_status = MagicMock()
        mock_get.return_value = mock_response

        result = extractors.extract_text_from_url("https://example.com/article")
        assert "main content" in result["text"].lower() or "article" in result["text"].lower()
        assert "source" in result
        assert result["source"] == "https://example.com/article"


def test_pdf_extraction_fallback(tmp_path):
    """Test PDF extraction with empty text (fallback to no OCR path)."""
    # Create a minimal PDF using fitz
    try:
        import fitz
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((50, 50), "PDF test content")
        pdf_path = tmp_path / "test.pdf"
        doc.save(pdf_path)
        doc.close()
    except ImportError:
        pytest.skip("fitz/pymupdf not available")
        return

    result = extractors.extract_text_from_pdf(pdf_path, use_ocr=False)
    assert "PDF test content" in result["text"]
    assert result["pages"] >= 1
    assert "source" in result


def test_pptx_extraction():
    """Test PPTX extraction with mock data."""
    try:
        from pptx import Presentation
    except ImportError:
        pytest.skip("python-pptx not available")
        return

    import tempfile
    with tempfile.TemporaryDirectory() as tmpdir:
        # Create a minimal PPTX
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        title.text = "Test Slide"
        
        pptx_path = Path(tmpdir) / "test.pptx"
        prs.save(pptx_path)

        # Extract
        result = extractors.extract_text_from_pptx(pptx_path)
        assert "Test Slide" in result["text"]
        assert result["slides"] >= 1


def test_embedding_adapter_local_provider():
    """Test embedding adapter with local provider."""
    class MockTokenizer:
        def encode(self, s, add_special_tokens=False):
            return [0] * max(1, len(s.split()))

    with patch("ingestion.embeddings.SentenceTransformer") as MockSentenceTransformer:
        mock_model = MagicMock()
        MockSentenceTransformer.return_value = mock_model
        
        # Mock encode to return simple arrays
        import numpy as np
        mock_model.encode.return_value = np.array([[0.1, 0.2], [0.3, 0.4]])

        adapter = embeddings.EmbeddingAdapter(model_name="test-model", provider="local")
        result = adapter.embed_texts(["text1", "text2"])
        
        assert len(result) == 2
        assert isinstance(result[0], list)
        assert len(result[0]) == 2


def test_embedding_adapter_openai_provider_missing_key():
    """Test that OpenAI provider fails gracefully without openai package or API key."""
    # Skip if openai is installed (test only relevant when it's not)
    try:
        import openai
        pytest.skip("openai package is installed; skipping test")
    except ImportError:
        pass
    
    with patch.dict("os.environ", {}, clear=True):
        try:
            adapter = embeddings.EmbeddingAdapter(model_name="text-embedding-3-small", provider="openai")
            assert False, "Should raise ImportError or ValueError"
        except (ImportError, ValueError) as e:
            # Either missing package or missing API key is acceptable
            assert "openai" in str(e).lower() or "api_key" in str(e).lower()


def test_chroma_isolation_by_user_notebook():
    """Test that Chroma collections isolate by user_id and notebook_id."""
    # Use EphemeralClient (in-memory) to avoid persistence/file locking issues on Windows
    store = vectorstore.ChromaAdapter(persist_directory=None)
    
    # Create collections for different users/notebooks
    col1 = store.get_or_create_collection("alice", "nb1")
    col2 = store.get_or_create_collection("alice", "nb2")
    col3 = store.get_or_create_collection("bob", "nb1")
    
    # Verify different names
    assert col1.name == "alice_nb1"
    assert col2.name == "alice_nb2"
    assert col3.name == "bob_nb1"
    
    # Upsert into col1
    chunks = [{"chunk_id": f"c{i}", "text": f"text{i}", "text_preview": "...", "source_id": "s1"} for i in range(2)]
    embeddings = [[0.1 * j for _ in range(10)] for j in range(len(chunks))]
    store.upsert_chunks("alice", "nb1", chunks, embeddings)
    
    # Verify col1 has data, others don't
    assert col1.count() == 2
    assert col2.count() == 0
    assert col3.count() == 0
