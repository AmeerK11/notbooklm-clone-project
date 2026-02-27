from __future__ import annotations

import io
import ipaddress
import socket
from pathlib import Path
from urllib.parse import urljoin, urlsplit, urlunsplit
from typing import Dict, Any

import requests
from bs4 import BeautifulSoup
from readability import Document
import fitz  # pymupdf
from pptx import Presentation


IPAddress = ipaddress.IPv4Address | ipaddress.IPv6Address

ALLOWED_URL_SCHEMES = {"http", "https"}
BLOCKED_LOCAL_HOST_SUFFIXES = (".localhost", ".local")
BLOCKED_EXACT_HOSTS = {"localhost", "localhost.localdomain"}
URL_FETCH_TIMEOUT = (5, 15)
URL_MAX_REDIRECTS = 3
URL_MAX_BYTES = 4 * 1024 * 1024
ALLOWED_URL_CONTENT_TYPES = ("text/html", "application/xhtml+xml", "text/plain")
REDIRECT_STATUS_CODES = {301, 302, 303, 307, 308}


class URLValidationError(ValueError):
    """Raised when a URL is unsafe or invalid for ingestion."""


def _is_disallowed_target_ip(ip_addr: IPAddress) -> bool:
    return (
        ip_addr.is_private
        or ip_addr.is_loopback
        or ip_addr.is_link_local
        or ip_addr.is_multicast
        or ip_addr.is_reserved
        or ip_addr.is_unspecified
    )


def _resolve_host_ips(hostname: str) -> set[IPAddress]:
    try:
        addr_info = socket.getaddrinfo(hostname, None, type=socket.SOCK_STREAM)
    except socket.gaierror as exc:
        raise URLValidationError(f"Unable to resolve host '{hostname}'.") from exc

    resolved_ips: set[IPAddress] = set()
    for _, _, _, _, sockaddr in addr_info:
        raw_ip = str(sockaddr[0]).split("%", 1)[0]
        try:
            resolved_ips.add(ipaddress.ip_address(raw_ip))
        except ValueError:
            continue

    if not resolved_ips:
        raise URLValidationError(f"Unable to resolve host '{hostname}'.")
    return resolved_ips


def _assert_public_hostname(hostname: str) -> None:
    hostname_clean = hostname.strip().rstrip(".").lower()
    if not hostname_clean:
        raise URLValidationError("URL must include a valid host.")
    if hostname_clean in BLOCKED_EXACT_HOSTS:
        raise URLValidationError("Local network URLs are not allowed.")
    if hostname_clean.endswith(BLOCKED_LOCAL_HOST_SUFFIXES):
        raise URLValidationError("Local network URLs are not allowed.")

    for ip_addr in _resolve_host_ips(hostname_clean):
        if _is_disallowed_target_ip(ip_addr):
            raise URLValidationError(
                f"Target host '{hostname_clean}' resolves to a restricted IP address."
            )


def validate_ingestion_url(url: str) -> str:
    candidate = (url or "").strip()
    if not candidate:
        raise URLValidationError("URL is required.")

    parsed = urlsplit(candidate)
    scheme = parsed.scheme.lower()
    if scheme not in ALLOWED_URL_SCHEMES:
        raise URLValidationError("Only http/https URLs are supported.")
    if not parsed.netloc:
        raise URLValidationError("URL must include a valid host.")
    if parsed.username or parsed.password:
        raise URLValidationError("URLs with embedded credentials are not allowed.")
    try:
        port = parsed.port
    except ValueError as exc:
        raise URLValidationError("URL includes an invalid port.") from exc
    if port is not None and not (1 <= port <= 65535):
        raise URLValidationError("URL includes an invalid port.")

    hostname = parsed.hostname or ""
    _assert_public_hostname(hostname)

    normalized_path = parsed.path or "/"
    return urlunsplit((scheme, parsed.netloc, normalized_path, parsed.query, parsed.fragment))


def _is_allowed_content_type(content_type: str) -> bool:
    kind = (content_type or "").split(";", 1)[0].strip().lower()
    if not kind:
        return True
    return any(kind.startswith(prefix) for prefix in ALLOWED_URL_CONTENT_TYPES)


def _read_response_text(response: requests.Response) -> str:
    content_type = response.headers.get("Content-Type", "")
    if not _is_allowed_content_type(content_type):
        raise URLValidationError(f"Unsupported content type for URL ingestion: '{content_type}'.")

    total_bytes = 0
    chunks: list[bytes] = []
    for chunk in response.iter_content(chunk_size=8192):
        if not chunk:
            continue
        total_bytes += len(chunk)
        if total_bytes > URL_MAX_BYTES:
            raise URLValidationError(
                f"URL response exceeded max size of {URL_MAX_BYTES} bytes."
            )
        chunks.append(chunk)

    payload = b"".join(chunks)
    encoding = response.encoding or response.apparent_encoding or "utf-8"
    return payload.decode(encoding, errors="replace")


def _fetch_html_from_url(url: str) -> tuple[str, str]:
    headers = {
        "User-Agent": "ingestion-bot/1.0",
        "Accept": "text/html,application/xhtml+xml,text/plain;q=0.9,*/*;q=0.1",
    }
    current_url = validate_ingestion_url(url)

    for hop in range(URL_MAX_REDIRECTS + 1):
        parsed = urlsplit(current_url)
        _assert_public_hostname(parsed.hostname or "")
        try:
            response = requests.get(
                current_url,
                timeout=URL_FETCH_TIMEOUT,
                headers=headers,
                allow_redirects=False,
                stream=True,
            )
        except requests.RequestException as exc:
            raise URLValidationError(f"Failed to fetch URL: {exc}") from exc

        try:
            if response.status_code in REDIRECT_STATUS_CODES:
                location = response.headers.get("Location", "").strip()
                if not location:
                    raise URLValidationError("Redirect response missing Location header.")
                if hop >= URL_MAX_REDIRECTS:
                    raise URLValidationError("Too many redirects while fetching URL.")
                current_url = validate_ingestion_url(urljoin(current_url, location))
                continue

            response.raise_for_status()
            return _read_response_text(response), current_url
        except requests.RequestException as exc:
            raise URLValidationError(f"Failed to fetch URL: {exc}") from exc
        finally:
            response.close()

    raise URLValidationError("Too many redirects while fetching URL.")


def extract_text_from_txt(path: Path) -> Dict[str, Any]:
    """Extract plain text from a .txt file.

    Returns a dict with keys: `text`, `pages` (1), `source` (path str).
    """
    text = None
    try:
        text = path.read_text(encoding="utf-8")
    except Exception:
        # fallback to latin1
        text = path.read_text(encoding="latin-1")
    return {"text": text, "pages": 1, "source": str(path)}


def extract_text_from_url(url: str) -> Dict[str, Any]:
    """Fetch a URL and extract the main article text using readability + BeautifulSoup.

    Returns dict with keys: `text`, `title`, `html`, `source`.
    """
    html, resolved_source_url = _fetch_html_from_url(url)
    doc = Document(html)
    summary_html = doc.summary()
    title = doc.title()
    soup = BeautifulSoup(summary_html, "html.parser")
    text = soup.get_text(separator="\n")
    return {"text": text, "title": title, "html": html, "source": resolved_source_url}


def extract_text_from_pdf(path: Path, use_ocr: bool = False) -> Dict[str, Any]:
    """Extract text from a PDF file using PyMuPDF (fitz).

    Args:
        path: Path to PDF file.
        use_ocr: If True, attempt OCR on images (requires pytesseract + tesseract).
                 Defaults to False (text-only extraction).

    Returns dict with keys: `text`, `pages`, `source`.
    """
    try:
        doc = fitz.open(path)
    except Exception as e:
        raise ValueError(f"Failed to open PDF {path}: {e}")

    text_parts = []
    page_count = len(doc)  # Capture page count before closing
    for page_num, page in enumerate(doc):
        try:
            text = page.get_text()
            if not text.strip() and use_ocr:
                # Attempt OCR on images in this page if text extraction is empty
                try:
                    import pytesseract
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                    img_data = pix.tobytes(output="png")
                    import PIL.Image
                    img = PIL.Image.open(io.BytesIO(img_data))
                    text = pytesseract.image_to_string(img)
                except Exception as ocr_err:
                    text = f"[OCR failed on page {page_num + 1}: {ocr_err}]"
            text_parts.append(text)
        except Exception as e:
            text_parts.append(f"[Error extracting page {page_num + 1}: {e}]")

    full_text = "\n".join(text_parts)
    doc.close()
    return {"text": full_text, "pages": page_count, "source": str(path)}


def extract_text_from_pptx(path: Path) -> Dict[str, Any]:
    """Extract text from a PowerPoint file (.pptx).

    Returns dict with keys: `text`, `slides`, `source`.
    """
    try:
        prs = Presentation(path)
    except Exception as e:
        raise ValueError(f"Failed to open PPTX {path}: {e}")

    slide_texts = []
    for slide_num, slide in enumerate(prs.slides):
        slide_text_parts = []
        for shape in slide.shapes:
            try:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text_parts.append(shape.text)
            except Exception:
                pass
        slide_texts.append("\n".join(slide_text_parts) if slide_text_parts else f"[Slide {slide_num + 1} - no text extracted]")

    full_text = "\n\n".join(slide_texts)
    return {"text": full_text, "slides": len(prs.slides), "source": str(path)}
